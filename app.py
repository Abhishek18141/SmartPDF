# PDF to PowerPoint Converter Flask App
# Optimized for Render Cloud Deployment
# Author: [Your Name]
# Version: 2.0

import os
import io
import re
import subprocess
from flask import Flask, jsonify, render_template, request, redirect, url_for, session, flash, send_file
from werkzeug.utils import secure_filename
import fitz  # PyMuPDF
from PIL import Image
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
import PyPDF2

# =============================================================================
# FLASK APP CONFIGURATION
# =============================================================================

app = Flask(__name__)
app.secret_key = os.environ.get("FLASK_SECRET_KEY", "your-secure-secret-key-here")

# Use /tmp for Render deployment (ephemeral storage)
UPLOAD_FOLDER = "/tmp/uploads"
TEMP_FOLDER = "/tmp/temp_uploads"

# Create directories
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(TEMP_FOLDER, exist_ok=True)

# =============================================================================
# AUTHENTICATION CONFIGURATION
# =============================================================================

USERNAME = "User"
PASSWORD = "Abc@123"

# =============================================================================
# UTILITY FUNCTIONS
# =============================================================================

def clear_temp_folder():
    """Clear temporary upload folder"""
    try:
        if os.path.exists(TEMP_FOLDER):
            for filename in os.listdir(TEMP_FOLDER):
                file_path = os.path.join(TEMP_FOLDER, filename)
                if os.path.isfile(file_path):
                    os.remove(file_path)
    except Exception:
        pass

def allowed_file(filename):
    """Check if file extension is allowed"""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() == 'pdf'

def clean_text(text, actual_page):
    """Clean extracted text from PDF"""
    lines = text.splitlines()
    cleaned = []
    
    for line in lines:
        l = line.strip()
        # Skip page number lines
        if re.fullmatch(rf"{actual_page}", l, re.IGNORECASE): 
            continue
        if re.fullmatch(rf"page\s*{actual_page}", l, re.IGNORECASE): 
            continue
        cleaned.append(l)
    
    joined = "\n".join(cleaned)
    joined = re.sub(r'\n{2,}', '<<<PARA>>>', joined)
    joined = re.sub(r'(?<!\n)\n(?!\n)', ' ', joined)
    return re.sub(r'\s{2,}', ' ', joined.replace('<<<PARA>>>', '\n\n')).strip()

def extract_pdf_data(pdf_path, start_page, end_page):
    """Extract text and images from PDF"""
    doc = fitz.open(pdf_path)
    total_pages = len(doc)

    # Validate range
    if total_pages == 0:
        raise ValueError("The PDF has no pages.")
    if not (1 <= start_page <= end_page <= total_pages):
        raise ValueError(f"Page range {start_page}-{end_page} is invalid. PDF has only {total_pages} page(s).")

    texts, images = [], []
    page_num = 1

    for i in range(start_page - 1, end_page):  # Convert to 0-based indexing
        try:
            page = doc[i]
            text = page.get_text()
            text = clean_text(text, i + 1)
            texts.append({'page_number': page_num, 'cleaned_text': text})

            page_images = []
            for img in page.get_images(full=True):
                base = doc.extract_image(img[0])
                image = Image.open(io.BytesIO(base['image']))
                page_images.append(image)
            if page_images:
                images.append({'page_number': page_num, 'images': page_images})

            page_num += 1
        except IndexError:
            raise IndexError(f"Page {i + 1} is out of bounds for this PDF.")

    doc.close()
    return texts, images

def generate_summary(texts, api_key):
    """Generate summaries using Gemini AI API"""
    summaries = []
    url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-flash:generateContent?key={api_key}"
    
    for page in texts:
        prompt = (
            "You are a helpful assistant creating PowerPoint slides.\n"
            "If the content is irrelevant or empty (e.g., 'Page 5'), respond with 'none'.\n"
            "If the content contains only a title (like 'Chapter 2'), return that title alone on the first line and write 'none' in the second line.\n"
            "If the content is useful, return:\n"
            "1. A short, meaningful title (max 6 words, just the title).\n"
            "2. 4–6 concise bullet points (each under 15 words).\n\n"
            f"Content:\n{page['cleaned_text']}"
        )

        try:
            response = requests.post(url, json={"contents": [{"parts": [{"text": prompt}]}]})

            if response.status_code != 200:
                # Extract and flash API error
                try:
                    error_msg = response.json().get('error', {}).get('message', 'Unknown error')
                except Exception:
                    error_msg = "Invalid response from Gemini API."

                flash(f"❌ Gemini API Error: {error_msg}", "danger")
                raise Exception(f"Gemini API Error: {error_msg}")
            
            result = response.json()['candidates'][0]['content']['parts'][0]['text']
            lines = result.strip().splitlines()
            
            if not lines:
                continue  # Skip empty response

            title = lines[0].strip()
            bullet_lines = [l.strip() for l in lines[1:] if l.strip()]
            bullet_text = "\n".join(bullet_lines)

            # If bullet says 'None' (title-only)
            if bullet_text.lower() == "none":
                summaries.append({
                    "page_number": page['page_number'],
                    "title": title,
                    "bullet_points": ""
                })
            # Normal title + bullets
            else:
                summaries.append({
                    "page_number": page['page_number'],
                    "title": title,
                    "bullet_points": bullet_text
                })
        except Exception as e:
            summaries.append({
                "page_number": page['page_number'],
                "title": f"Page {page['page_number']}",
                "bullet_points": "Error occurred."
            })
    return summaries

def build_ppt(bullet_points, images, output_path):
    """Build PowerPoint presentation"""
    prs = Presentation()
    layout = prs.slide_layouts[6]  # blank layout

    b_map = {b['page_number']: b for b in bullet_points}
    i_map = {i['page_number']: i['images'] for i in images}

    for pg in sorted(b_map):
        bp_data = b_map[pg]
        title = bp_data['title'].strip()
        title_lower = title.lower()
        bullets = [line.strip() for line in bp_data['bullet_points'].split('\n') if line.strip() and line.lower() != "none"]
        imgs = i_map.get(pg, [])

        # Skip if no title, no bullets, and no images
        if title_lower == "none" and not bullets and not imgs:
            continue

        slide = prs.slides.add_slide(layout)
        sw, sh = prs.slide_width, prs.slide_height
        margin = Inches(0.3)
        top = Inches(0.6)

        # Case 1: Only title
        if title_lower != "none" and not bullets and not imgs:
            title_box = slide.shapes.add_textbox(margin, sh / 2 - Inches(0.5), sw - 2 * margin, Inches(1))
            tf = title_box.text_frame
            tf.text = title
            tf.paragraphs[0].font.size = Pt(32)
            tf.paragraphs[0].alignment = 1  # Center
            continue  # no more content to add

        # Case 2: Only images (no title or bullets)
        if title_lower == "none" and not bullets and imgs:
            n = len(imgs)
            if n == 1:
                img = imgs[0].convert("RGB")
                temp = os.path.join(TEMP_FOLDER, f"temp_{pg}.png")
                img.save(temp)
                slide.shapes.add_picture(temp, 0, 0, width=sw, height=sh)
                os.remove(temp)
            else:
                cols = min(n, 4)
                rows = (n + cols - 1) // cols
                iw = (sw - margin * (cols + 1)) / cols
                ih = (sh - margin * (rows + 1)) / rows
                for i, img in enumerate(imgs):
                    temp = os.path.join(TEMP_FOLDER, f"temp_{pg}_{i}.png")
                    img.convert("RGB").save(temp)
                    c, r = i % cols, i // cols
                    l = margin + c * (iw + margin)
                    t = margin + r * (ih + margin)
                    slide.shapes.add_picture(temp, l, t, width=iw, height=ih)
                    os.remove(temp)
            continue  # done with this slide

        # Case 3: Title and/or images and/or bullets
        if title_lower != "none":
            title_box = slide.shapes.add_textbox(margin, Inches(0.1), sw - 2 * margin, Inches(0.5))
            tf = title_box.text_frame
            tf.text = title
            tf.paragraphs[0].font.size = Pt(24)
            tf.paragraphs[0].alignment = 1
            top = Inches(0.6)

        # Images (top half)
        if imgs:
            n = len(imgs)
            if n == 1:
                img = imgs[0].convert("RGB")
                temp = os.path.join(TEMP_FOLDER, f"temp_{pg}.png")
                img.save(temp)
                iw, ih = sw / 4, sh / 4
                slide.shapes.add_picture(temp, margin, top, width=iw, height=ih)
                slide.shapes.add_picture(temp, sw - iw - margin, top, width=iw, height=ih)
                os.remove(temp)
                top += ih + Inches(0.5)
            elif n > 1:
                cols = min(n, 4)
                rows = (n + cols - 1) // cols
                iw = (sw - margin * (cols + 1)) / cols
                ih = (sh / 2 - margin * (rows + 1)) / rows
                for i, img in enumerate(imgs):
                    temp = os.path.join(TEMP_FOLDER, f"temp_{pg}_{i}.png")
                    img.convert("RGB").save(temp)
                    c, r = i % cols, i // cols
                    l = margin + c * (iw + margin)
                    t = top + r * (ih + margin)
                    slide.shapes.add_picture(temp, l, t, width=iw, height=ih)
                    os.remove(temp)
                top += ih * rows + Inches(0.3)

        # Bullet points
        if bullets:
            box = slide.shapes.add_textbox(margin, top, sw - 2 * margin, sh - top - margin)
            tf = box.text_frame
            for line in bullets:
                p = tf.add_paragraph()
                p.text = f"• {line.strip('* ').strip()}"
                p.level = 0
                p.font.size = Pt(16)

    prs.save(output_path)

# =============================================================================
# FLASK ROUTES
# =============================================================================

@app.route('/')
def index():
    """Home route - redirect to login"""
    return redirect(url_for('login'))

@app.route('/login', methods=['GET', 'POST'])
def login():
    """User authentication"""
    clear_temp_folder()
    session.clear()
    
    if request.method == 'POST':
        if request.form['username'] == USERNAME and request.form['password'] == PASSWORD:
            session['logged_in'] = True
            return redirect(url_for('apikey'))
        flash('Invalid credentials')
    
    return render_template('login.html')

@app.route('/about')
def about():
    """About page"""
    return render_template('about.html')

@app.route('/logout')
def logout():
    """User logout"""
    clear_temp_folder()
    session.clear()
    flash("You have been logged out.")
    return redirect(url_for('login'))

@app.route('/apikey', methods=['GET', 'POST'])
def apikey():
    """API key configuration"""
    if not session.get('logged_in'):
        return redirect(url_for('login'))
    
    if request.method == 'POST':
        session['api_key'] = request.form['api_key']
        return redirect(url_for('upload'))
    
    return render_template('apikey.html')

@app.route('/get_page_count', methods=['POST'])
def get_page_count():
    """Get total pages from uploaded PDF"""
    file = request.files.get('pdf_file')
    if not file or not file.filename.lower().endswith('.pdf'):
        return jsonify({'success': False, 'error': 'Invalid file. Please upload a PDF file.'})

    try:
        os.makedirs(TEMP_FOLDER, exist_ok=True)
        filepath = os.path.join(TEMP_FOLDER, secure_filename(file.filename))
        file.save(filepath)

        reader = PyPDF2.PdfReader(filepath)
        total_pages = len(reader.pages)

        return jsonify({'success': True, 'total_pages': total_pages})
    except Exception as e:
        return jsonify({'success': False, 'error': f'Error reading PDF: {str(e)}'})

@app.route('/upload', methods=['GET', 'POST'])
def upload():
    """File upload and processing configuration"""
    if request.method == 'POST':
        pdf_file = request.files.get('pdf_file')
        if not pdf_file or not allowed_file(pdf_file.filename):
            flash("Invalid file format. Please upload a PDF.", "danger")
            return redirect(request.url)

        filename = secure_filename(pdf_file.filename)
        filepath = os.path.join(TEMP_FOLDER, filename)

        try:
            start_page = int(request.form.get('start_page'))
            end_page = int(request.form.get('end_page'))
        except (ValueError, TypeError):
            flash("Invalid page numbers.", "danger")
            return redirect(request.url)

        pdf_file.save(filepath)

        with open(filepath, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            total_pages = len(pdf_reader.pages)

        if start_page < 1 or end_page < 1 or start_page > end_page or end_page > total_pages:
            flash("Invalid page range. Check your input.", "danger")
            return redirect(request.url)

        session['pdf_path'] = filepath
        session['start_page'] = start_page
        session['end_page'] = end_page

        action = request.form.get('action')
        if action == 'ppt':
            return redirect(url_for('ppt'))
        elif action == 'qa':
            return redirect(url_for('qa'))
        else:
            flash("Unknown action selected.", "danger")
            return redirect(request.url)

    return render_template('upload.html')

@app.route('/qa', methods=['GET', 'POST'])
def qa():
    """Question and Answer functionality"""
    if not session.get('pdf_path'): 
        return redirect(url_for('upload'))
    
    texts, _ = extract_pdf_data(session['pdf_path'], session['start_page'], session['end_page'])
    combined = "\n\n".join(p['cleaned_text'] for p in texts)
    answer = ""
    
    if request.method == 'POST':
        question = request.form['question']
        prompt = f"Answer the question based on this text:\n{combined}\n\nQuestion: {question}"
        api_key = session.get('api_key')
        url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-flash:generateContent?key={api_key}"
        
        try:
            res = requests.post(url, json={"contents": [{"parts": [{"text": prompt}]}]})
            res.raise_for_status()
            answer = res.json()['candidates'][0]['content']['parts'][0]['text']
        except Exception as e:
            answer = f"Error: {e}"
    
    return render_template('qa.html', answer=answer)

@app.route('/ppt', methods=['GET', 'POST'])
def ppt():
    """PowerPoint generation"""
    pdf_path = session.get('pdf_path')
    start_page = session.get('start_page')
    end_page = session.get('end_page')
    api_key = session.get('api_key')

    if not pdf_path or start_page is None or end_page is None:
        flash("Session data missing. Please upload the PDF again.", "danger")
        return redirect(url_for('upload'))

    if request.method == 'POST':
        filename = request.form.get('filename')
        if not filename:
            flash("Please enter a valid PPT filename.", "danger")
            return render_template('ppt.html')

        try:
            # Extract text and images
            texts, images = extract_pdf_data(pdf_path, start_page, end_page)

            # Generate summaries from texts (calls Gemini API)
            summaries = generate_summary(texts, api_key)

            if not summaries:
                return render_template('ppt.html', generated=True, no_content=True)

            # Build the PPT file path
            ppt_filename = f"{filename}.pptx"
            ppt_path = os.path.join(TEMP_FOLDER, ppt_filename)

            # Create PPT
            build_ppt(summaries, images, ppt_path)

            # Save ppt path in session for download
            session['ppt_filename'] = ppt_path

            # Pass to template to show download link and summary
            return render_template('ppt.html', generated=True, filename=filename, summaries=summaries, no_content=False)

        except Exception as e:
            flash(f"Error while generating PPT: {str(e)}", "danger")
            return render_template('ppt.html')

    # GET request, just render the page normally
    return render_template('ppt.html')

@app.route('/download_ppt')
def download_ppt():
    """Download generated PowerPoint file"""
    filename = session.get('ppt_filename', 'output_presentation.pptx')
    return send_file(filename, as_attachment=True)

# =============================================================================
# APPLICATION STARTUP
# =============================================================================

if __name__ == '__main__':
    # Get port from environment variable for Render deployment
    port = int(os.environ.get('PORT', 5000))
    app.run(host='0.0.0.0', port=port, debug=False)
